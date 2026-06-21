"""
Generate the SR Inc. executive briefing deck:
  "How the ERCOT Settlement Forecast Works"

Run:
  .venv/bin/python build_forecast_deck.py
Output:
  SR_ERCOT_Forecast_Methodology.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── SR Inc. palette ───────────────────────────────────────────────────────────
SR_BLUE       = RGBColor(0x00, 0x69, 0xB3)
SR_BLUE_DARK  = RGBColor(0x00, 0x55, 0x8E)
SR_BLUE_LIGHT = RGBColor(0x54, 0xA4, 0xDA)
SR_BLUE_PALE  = RGBColor(0xD7, 0xE2, 0xF2)
SR_BLUE_GHOST = RGBColor(0xEC, 0xF0, 0xF9)
SR_GREEN      = RGBColor(0x88, 0xA9, 0x18)
SR_GREEN_ALT  = RGBColor(0x01, 0xA0, 0x6F)
SR_DARK_GREY  = RGBColor(0x63, 0x66, 0x6F)
SR_MID_GREY   = RGBColor(0x84, 0x84, 0x84)
SR_NAVY       = RGBColor(0x00, 0x1F, 0x5F)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BAD_RED       = RGBColor(0xB2, 0x3A, 0x48)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank

# ── helpers ───────────────────────────────────────────────────────────────────

def rgb_hex(r):
    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_pt=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=18, bold=False, italic=False,
             color=SR_DARK_GREY, align=PP_ALIGN.LEFT,
             wrap=True, bg_rgb=None):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, font_name="Calibri", font_size=16, bold=False,
             color=SR_DARK_GREY, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def slide_header(slide, eyebrow, title, subtitle=""):
    """SR-branded slide header: blue top band, green eyebrow, white title."""
    add_rect(slide, 0, 0, 13.33, 1.55, fill_rgb=SR_BLUE)
    add_rect(slide, 0.45, 0.18, 0.6, 0.06, fill_rgb=SR_GREEN)  # green rule
    add_text(slide, eyebrow.upper(), 0.45, 0.26, 9, 0.3,
             font_size=9, bold=True, color=RGBColor(0xCF,0xE6,0xB0),
             font_name="Calibri")
    add_text(slide, title, 0.45, 0.48, 11, 0.75,
             font_size=28, bold=True, color=WHITE, font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle, 0.45, 1.18, 11, 0.35,
                 font_size=13, color=RGBColor(0xCC,0xDD,0xEE), font_name="Calibri")

def bullet_box(slide, left, top, width, height, title, bullets,
               title_color=SR_BLUE, bullet_color=SR_DARK_GREY,
               bg=SR_BLUE_GHOST, border=SR_BLUE_PALE,
               title_size=15, bullet_size=13):
    add_rect(slide, left, top, width, height,
             fill_rgb=bg, line_rgb=border, line_pt=0.75)
    # title
    add_text(slide, title, left+0.18, top+0.13, width-0.3, 0.3,
             font_size=title_size, bold=True, color=title_color, font_name="Calibri")
    # bullets
    y = top + 0.48
    line_h = (height - 0.55) / max(len(bullets), 1)
    for b in bullets:
        add_text(slide, f"• {b}", left+0.22, y, width-0.4, line_h+0.05,
                 font_size=bullet_size, color=bullet_color, font_name="Calibri", wrap=True)
        y += line_h

def footer_bar(slide, text="Sustainability Roundtable, Inc.  ·  Confidential"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=SR_BLUE_GHOST)
    add_text(slide, text, 0.3, 7.17, 12, 0.28,
             font_size=9, color=SR_MID_GREY, font_name="Calibri")
    add_text(slide, "sr-inc.com", 11.5, 7.17, 1.6, 0.28,
             font_size=9, color=SR_BLUE, font_name="Calibri", align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Full blue background
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=SR_NAVY)
# Accent stripe
add_rect(s, 0, 6.8, 13.33, 0.7, fill_rgb=SR_BLUE)
# Green rule
add_rect(s, 0.8, 2.5, 0.9, 0.07, fill_rgb=SR_GREEN)

add_text(s, "SUSTAINABILITY ROUNDTABLE, INC.", 0.8, 1.6, 11, 0.4,
         font_size=10, bold=True, color=RGBColor(0xCF,0xE6,0xB0), font_name="Calibri")
add_text(s, "How the ERCOT\nSettlement Forecast Works", 0.8, 2.65, 11.5, 1.8,
         font_size=40, bold=True, color=WHITE, font_name="Calibri")
add_text(s, "A plain-language guide for executive stakeholders", 0.8, 4.55, 11, 0.5,
         font_size=18, color=SR_BLUE_LIGHT, font_name="Calibri")
add_text(s, "June 2026  ·  ERCOT Customer Portals  ·  Confidential", 0.8, 5.2, 11, 0.4,
         font_size=13, color=SR_MID_GREY, font_name="Calibri")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What the Forecast Is (and Isn't)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Overview", "What the Forecast Is — and Isn't",
             subtitle="Setting the context before the methodology")
footer_bar(s)

add_text(s, "Purpose", 0.5, 1.75, 5.8, 0.35, font_size=13, bold=True,
         color=SR_GREEN, font_name="Calibri")
add_text(s,
    "The Projected Bill tab gives offtakers a rolling estimate of the current month's "
    "and next month's net settlement under their VPPA or CfD contract — before ERCOT "
    "publishes final settlement data (which lags ~60 days).",
    0.5, 2.1, 5.8, 1.0, font_size=13, color=SR_DARK_GREY, font_name="Calibri", wrap=True)

bullet_box(s, 0.5, 3.2, 5.8, 2.7, "What it IS",
    ["Weather-driven generation model calibrated to the plant's real metered output",
     "Forward price assumption entered by the user (or defaulted to trailing capture)",
     "Transparent: every assumption is visible and overridable",
     "Updated automatically each time the page loads"],
    bg=SR_BLUE_GHOST, title_color=SR_GREEN_ALT)

bullet_box(s, 6.9, 1.75, 6.0, 4.15, "What it is NOT",
    ["Final or invoiced settlement — that appears on the Past Settlement page",
     "A price forecast — the user supplies the forward price assumption",
     "A guarantee of future generation — weather and curtailment can differ",
     "An audited figure — for billing disputes, use the ERCOT data on the Invoice Audit page",
     "Retroactively binding — the model is refreshed on each visit; past estimates are not stored"],
    bg=RGBColor(0xFD,0xF0,0xF2), border=RGBColor(0xE8,0xC0,0xC6),
    title_color=BAD_RED, bullet_color=SR_DARK_GREY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — End-to-End Flow
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Architecture", "End-to-End: From Weather to Dollar",
             subtitle="Five steps, three data sources, one settlement number")
footer_bar(s)

steps = [
    ("1", "Weather\nForecast", "Open-Meteo API\n16-day hi-res\n+35-day GEFS\n+ERA5 clim.", SR_BLUE),
    ("2", "Height\nExtrap.", "Wind speed\nscaled to\nhub height\n(Hellmann 1/7)", SR_BLUE_LIGHT),
    ("3", "Generation\nModel", "Solar: GHI÷1000\n× capacity MW\nWind: cubic\npower curve", SR_GREEN),
    ("4", "Calibration\nFactor", "Model ÷ SCED\nover last 60 days\nAnchors to\nreal output", SR_GREEN_ALT),
    ("5", "Settlement\nCalc.", "Daily MWh ×\n(market price\n− strike)\n= net $", SR_NAVY),
]
box_w = 2.15
for i, (num, title, detail, col) in enumerate(steps):
    x = 0.4 + i * (box_w + 0.18)
    add_rect(s, x, 1.65, box_w, 0.55, fill_rgb=col)
    add_text(s, f"STEP {num}", x+0.08, 1.68, box_w-0.1, 0.28,
             font_size=8, bold=True, color=RGBColor(0xCC,0xEE,0xFF), font_name="Calibri")
    add_text(s, title, x+0.08, 1.92, box_w-0.1, 0.28,
             font_size=13, bold=True, color=WHITE, font_name="Calibri")
    add_rect(s, x, 2.2, box_w, 2.6, fill_rgb=SR_BLUE_GHOST,
             line_rgb=SR_BLUE_PALE, line_pt=0.75)
    add_text(s, detail, x+0.12, 2.3, box_w-0.2, 2.4,
             font_size=12, color=SR_DARK_GREY, font_name="Calibri", wrap=True)
    if i < 4:
        add_text(s, "→", x+box_w+0.01, 2.85, 0.22, 0.4,
                 font_size=20, bold=True, color=SR_BLUE, font_name="Calibri")

add_text(s, "Output: daily net settlement bar chart + month total, refreshed on every page load.",
         0.5, 5.05, 12, 0.4, font_size=12, color=SR_MID_GREY,
         font_name="Calibri", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Weather Data Three-Tier Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Data Sources", "Three-Tier Weather Architecture",
             subtitle="High-res near-term → ensemble mid-range → climatological tail")
footer_bar(s)

tiers = [
    ("Tier 1 — Standard Forecast", "Days 1–16",
     ["Source: Open-Meteo (free, no API key)",
      "Model: ECMWF / GFS high-resolution NWP",
      "Solar: hourly GHI (W/m²), temperature",
      "Wind: hourly speed at 80 m and 120 m",
      "Updated: every page load (2-hr file cache)",
      "Also fetches 31 past days for context"],
     SR_BLUE),
    ("Tier 2 — GEFS Ensemble P50", "Days 17–35",
     ["Source: Open-Meteo Ensemble API",
      "Model: GFS05 — 31 members averaged",
      "Ensemble mean = P50 (expected value)",
      "Wind: 80 m extrapolated to 120 m",
      "Updated: every 6 hours (4×/day)",
      "Boundary artifact trimmed automatically"],
     SR_GREEN),
    ("Tier 3 — Prior-Year ERA5", "Days 36–Month End",
     ["Source: Open-Meteo Archive API (ERA5)",
      "Same calendar days from prior year",
      "Realistic day-to-day weather variation",
      "Better than a flat monthly average",
      "Cached 24 hrs (historical, stable)",
      "Flat historical shape used as backstop"],
     SR_GREEN_ALT),
]

for i, (title, period, bullets, col) in enumerate(tiers):
    x = 0.4 + i * 4.28
    add_rect(s, x, 1.65, 4.05, 0.6, fill_rgb=col)
    add_text(s, period, x+0.15, 1.68, 3.7, 0.25,
             font_size=10, bold=True, color=RGBColor(0xCC,0xEE,0xFF), font_name="Calibri")
    add_text(s, title, x+0.15, 1.88, 3.7, 0.35,
             font_size=13, bold=True, color=WHITE, font_name="Calibri")
    add_rect(s, x, 2.25, 4.05, 3.6, fill_rgb=SR_BLUE_GHOST,
             line_rgb=SR_BLUE_PALE, line_pt=0.75)
    y = 2.38
    for b in bullets:
        add_text(s, f"• {b}", x+0.18, y, 3.7, 0.38,
                 font_size=12, color=SR_DARK_GREY, font_name="Calibri", wrap=True)
        y += 0.53

add_text(s, "Source priority: settled SCED actuals  >  Tier 1  >  Tier 2  >  Tier 3  >  flat historical shape",
         0.5, 6.05, 12.3, 0.4, font_size=11, color=SR_NAVY, bold=True,
         font_name="Calibri", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Generation Models
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Physics Models", "How Weather Becomes MWh",
             subtitle="Separate models for solar PV and wind — both anchored by a calibration factor")
footer_bar(s)

# Solar box
add_rect(s, 0.4, 1.65, 5.9, 4.7, fill_rgb=SR_BLUE_GHOST, line_rgb=SR_BLUE_PALE, line_pt=0.75)
add_rect(s, 0.4, 1.65, 5.9, 0.5, fill_rgb=SR_BLUE)
add_text(s, "☀  Solar PV Model", 0.58, 1.68, 5.5, 0.4,
         font_size=15, bold=True, color=WHITE, font_name="Calibri")

solar_items = [
    ("Formula", "MW = Capacity × (GHI W/m² ÷ 1,000) × Cal. Factor"),
    ("Input", "Hourly shortwave (GHI) from weather API"),
    ("Capacity", "Contracted MW share of the plant nameplate"),
    ("Cal. Factor", "Absorbs: panel efficiency (~17–20%), DC/AC losses,\n"
                    "soiling, row shading, availability derating"),
    ("Typical range", "Cal. Factor ≈ 0.15–0.25 for a real PV plant\n"
                      "(model is intentionally raw; factor corrects it)"),
    ("Cap", "Output capped at contracted capacity MW"),
]
y = 2.32
for label, detail in solar_items:
    add_text(s, label, 0.58, y, 1.5, 0.38, font_size=11, bold=True,
             color=SR_BLUE, font_name="Calibri")
    add_text(s, detail, 2.1, y, 4.0, 0.42, font_size=11,
             color=SR_DARK_GREY, font_name="Calibri", wrap=True)
    y += 0.6

# Wind box
add_rect(s, 6.9, 1.65, 6.0, 4.7, fill_rgb=SR_BLUE_GHOST, line_rgb=SR_BLUE_PALE, line_pt=0.75)
add_rect(s, 6.9, 1.65, 6.0, 0.5, fill_rgb=SR_GREEN)
add_text(s, "🌬  Wind Model", 7.08, 1.68, 5.7, 0.4,
         font_size=15, bold=True, color=WHITE, font_name="Calibri")

wind_items = [
    ("Step 1", "Extrapolate wind speed from 80/120 m\nto turbine hub height (Hellmann α = 1/7)"),
    ("Step 2", "Apply turbine power curve:\n"
               "  < cut-in  →  0 MW\n"
               "  cut-in → rated  →  cubic ramp\n"
               "  > rated  →  full capacity\n"
               "  > cut-out  →  0 MW (storm protection)"),
    ("Turbine\nparams", "Per asset: Vestas V110 cut-out 20 m/s;\nNordex N149 cut-out 25 m/s"),
    ("Cal. Factor", "Scales model to real SCED output;\ntypically 0.5–1.5 for wind"),
]
y = 2.32
for label, detail in wind_items:
    add_text(s, label, 7.08, y, 1.55, 0.48, font_size=11, bold=True,
             color=SR_GREEN, font_name="Calibri")
    add_text(s, detail, 8.65, y, 4.05, 0.6, font_size=11,
             color=SR_DARK_GREY, font_name="Calibri", wrap=True)
    y += 0.78

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Calibration
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Calibration", "Anchoring the Model to Real Plant Output",
             subtitle="The calibration factor is the single most important input")
footer_bar(s)

bullet_box(s, 0.4, 1.65, 5.9, 2.7, "What is the Calibration Factor?",
    ["A scalar multiplier applied to model output so the\n"
     "long-run average matches the plant's real SCED metering",
     "Computed as:  Actual MWh  ÷  Model MWh\n"
     "over the most recent 60 days of settled SCED data",
     "Absorbs all systematic differences between the simple\n"
     "weather model and the real plant (efficiency, losses,\n"
     "availability, curtailment, measurement error)",
     "Clipped to [0.05, 3.0] — values outside this range\n"
     "indicate a data problem, not a real plant behaviour"],
    title_color=SR_BLUE, bg=SR_BLUE_GHOST)

bullet_box(s, 0.4, 4.55, 5.9, 1.85, "Data Source for Calibration",
    ["Uses ERA5 archive (not the live forecast) for weather —\n"
     "live forecast zeros out radiation beyond ~30 days,\n"
     "but SCED lags ~60 days, so archive is essential",
     "SCED data sourced from ERCOT 15-min generation tables\n"
     "and scaled to the offtaker's contracted volume share"],
    title_color=SR_BLUE, bg=SR_BLUE_GHOST)

bullet_box(s, 6.9, 1.65, 6.0, 2.7, "Example Cal. Factors (Live)",
    ["Markum Solar:  1.231  (from 61 SCED days)\n"
     "  → model produces 82% of actual output;\n"
     "    factor corrects upward by 23%",
     "Azure Sky Wind:  1.108  (from 60+ SCED days)\n"
     "  → model slightly understates output;\n"
     "    factor adds ~11%",
     "Hidalgo Wind:  derived from available SCED history"],
    title_color=SR_GREEN_ALT, bg=SR_BLUE_GHOST)

bullet_box(s, 6.9, 4.55, 6.0, 1.85, "Consistency Guarantee",
    ["The calibration step and the forecast step always\n"
     "use the same power curve parameters (cut-in, rated\n"
     "speed, cut-out) — ensuring the factor is meaningful\n"
     "and not distorted by a model mismatch"],
    title_color=SR_GREEN_ALT, bg=SR_BLUE_GHOST)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Settlement Calculation
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Settlement Math", "From MWh to Dollars",
             subtitle="VPPA / CfD settlement logic applied daily")
footer_bar(s)

# Formula box
add_rect(s, 0.5, 1.65, 12.3, 1.1, fill_rgb=SR_NAVY)
add_text(s, "Daily Net Settlement ($)  =  MWh  ×  (Market Price  −  Strike Price)",
         0.7, 1.72, 12.0, 0.6, font_size=22, bold=True,
         color=WHITE, font_name="Calibri", align=PP_ALIGN.CENTER)
add_text(s, "Sign convention: positive = offtaker receives;  negative = offtaker pays",
         0.7, 2.22, 12.0, 0.4, font_size=12,
         color=SR_BLUE_LIGHT, font_name="Calibri", align=PP_ALIGN.CENTER)

bullet_box(s, 0.5, 2.9, 4.0, 3.7, "MWh",
    ["Forecasted daily generation at\nthe contracted volume share",
     "Source priority:\n  1. Settled SCED (actuals)\n  2. Weather model × cal. factor\n  3. GEFS ensemble\n  4. ERA5 climatology",
     "Capped at contracted capacity MW"],
    title_color=SR_BLUE)

bullet_box(s, 4.7, 2.9, 4.0, 3.7, "Market Price",
    ["Real-time settlement point price\n($/MWh) at the contract node",
     "For past settled days: actual\nERCOT 15-min RT prices",
     "For forecast days: user-entered\nforward price assumption\n(defaults to trailing capture\nprice from history)",
     "Not a modelled price — always\na user assumption"],
    title_color=SR_GREEN)

bullet_box(s, 8.9, 2.9, 4.0, 3.7, "Strike Price",
    ["Fixed contract price in $/MWh\n(set in Contract Terms page)",
     "Markum Solar:  $35.00/MWh",
     "Azure Sky Wind:  $17.34/MWh",
     "Hidalgo Wind:  $35.00/MWh",
     "When market > strike:\n  offtaker RECEIVES the spread",
     "When market < strike:\n  offtaker PAYS the spread"],
    title_color=SR_BLUE_LIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What Drives the Numbers
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Key Drivers", "What Moves the Forecast",
             subtitle="Ranked by typical impact on the monthly net settlement")
footer_bar(s)

drivers = [
    ("1", "Forward Price\nAssumption", "Largest single driver. The user sets this on the sidebar. "
      "A $5/MWh change on a 40,000 MWh month shifts the estimate by ±$200,000. "
      "Default is the trailing capture price from history.",
      SR_BLUE, "High"),
    ("2", "Generation\nForecast", "Weather uncertainty compounds over the month. Near-term (16 days) "
      "is high-res and reliable. GEFS ensemble is P50 and less certain. "
      "ERA5 climatology adds typical-year variation.",
      SR_GREEN, "Medium–High"),
    ("3", "Calibration\nFactor", "Anchors the model to real output. Estimated from 60 days of SCED. "
      "Stable once enough history is available; unreliable on new plants "
      "with < 30 days of settled data.",
      SR_GREEN_ALT, "Medium"),
    ("4", "Market Price\nVariability", "Actual ERCOT prices are highly volatile (weather, demand, outages). "
      "The forward assumption smooths this. Large deviation events (e.g. winter storms) "
      "are not captured.",
      SR_BLUE_LIGHT, "Low (estimate)"),
]

for i, (num, title, detail, col, impact) in enumerate(drivers):
    y = 1.7 + i * 1.3
    add_rect(s, 0.4, y, 0.55, 1.1, fill_rgb=col)
    add_text(s, num, 0.4, y+0.3, 0.55, 0.5,
             font_size=22, bold=True, color=WHITE,
             font_name="Calibri", align=PP_ALIGN.CENTER)
    add_rect(s, 0.95, y, 2.2, 1.1, fill_rgb=SR_BLUE_GHOST, line_rgb=SR_BLUE_PALE, line_pt=0.5)
    add_text(s, title, 1.1, y+0.2, 1.9, 0.7,
             font_size=13, bold=True, color=col, font_name="Calibri")
    add_rect(s, 3.15, y, 8.8, 1.1, fill_rgb=SR_BLUE_GHOST, line_rgb=SR_BLUE_PALE, line_pt=0.5)
    add_text(s, detail, 3.3, y+0.12, 8.4, 0.88,
             font_size=12, color=SR_DARK_GREY, font_name="Calibri", wrap=True)
    add_rect(s, 11.95, y, 1.0, 1.1, fill_rgb=col)
    add_text(s, impact, 11.97, y+0.28, 0.95, 0.55,
             font_size=10, bold=True, color=WHITE,
             font_name="Calibri", align=PP_ALIGN.CENTER, wrap=True)

add_text(s, "← Impact", 12.0, 1.55, 1.2, 0.28,
         font_size=9, color=SR_MID_GREY, font_name="Calibri")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Limitations & Transparency
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Methodology Notes", "Limitations & Transparency",
             subtitle="What the model cannot do — and where to find the authoritative figures")
footer_bar(s)

bullet_box(s, 0.4, 1.65, 5.9, 4.7, "Known Limitations",
    ["Weather beyond 35 days uses prior-year climatology —\n"
     "actual weather can differ significantly",
     "The power curve is a simplified cubic model;\n"
     "real turbine curves have smoother transitions",
     "No curtailment model — forecasted output may exceed\n"
     "what ERCOT actually dispatches",
     "Forward price is user-supplied — the tool does not\n"
     "forecast ERCOT market prices",
     "Calibration requires ~30+ days of settled SCED data;\n"
     "results on new plants are less reliable",
     "Extreme weather events (storms, heat waves) are not\n"
     "captured in the ensemble P50"],
    title_color=BAD_RED, bg=RGBColor(0xFD,0xF0,0xF2),
    border=RGBColor(0xE8,0xC0,0xC6))

bullet_box(s, 6.9, 1.65, 6.0, 2.25, "Authoritative Data Sources",
    ["Past Settlement (actual): ERCOT 15-min metered\n"
     "generation × real-time settlement prices",
     "Invoice Audit page: full interval-level detail,\n"
     "downloadable for billing verification",
     "Cross-check: EIA Form 923 monthly generation\n"
     "shown alongside SCED for independent validation"],
    title_color=SR_GREEN_ALT)

bullet_box(s, 6.9, 4.1, 6.0, 2.25, "Free & Cited Data Sources",
    ["Weather: Open-Meteo (open source, no API key)\n"
     "  open-meteo.com",
     "Reanalysis: ERA5 via Open-Meteo Archive API\n"
     "  (ECMWF reanalysis v5, hourly, global)",
     "Ensemble: GFS05 via Open-Meteo Ensemble API\n"
     "  (NOAA GFS, 31-member, 35-day horizon)",
     "Generation: ERCOT SCED / MIS public data\n"
     "  (15-min, plant-level, ~60-day lag)"],
    title_color=SR_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary / Thank You
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=SR_NAVY)
add_rect(s, 0, 6.5, 13.33, 1.0, fill_rgb=SR_BLUE)
add_rect(s, 0.8, 3.55, 1.1, 0.07, fill_rgb=SR_GREEN)

add_text(s, "SUSTAINABILITY ROUNDTABLE, INC.", 0.8, 1.8, 11, 0.4,
         font_size=10, bold=True, color=RGBColor(0xCF,0xE6,0xB0), font_name="Calibri")
add_text(s, "Key Takeaways", 0.8, 2.25, 11, 0.65,
         font_size=34, bold=True, color=WHITE, font_name="Calibri")

takeaways = [
    "The forecast uses real weather data updated daily — not a static model",
    "The calibration factor is the critical link between the weather model and real plant output",
    "Near-term (16 days) is reliable; beyond 35 days is climatological context",
    "The forward price assumption drives the dollar outcome — always user-controlled",
    "Authoritative settlement figures are on the Past Settlement page, not the forecast",
]
y = 3.7
for t in takeaways:
    add_rect(s, 0.75, y, 0.12, 0.32, fill_rgb=SR_GREEN)
    add_text(s, t, 1.05, y, 11, 0.38, font_size=14,
             color=WHITE, font_name="Calibri", wrap=True)
    y += 0.52

add_text(s, "Questions?  Contact Sustainability Roundtable, Inc.  ·  sr-inc.com",
         0.8, 6.58, 11, 0.35, font_size=11, color=SR_BLUE_LIGHT, font_name="Calibri")

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/michaelbarry/Documents/Github/ercot-suite/SR_ERCOT_Forecast_Methodology.pptx"
prs.save(out)
print(f"Saved: {out}")
